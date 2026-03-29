# import os
# import bcrypt
# import pymongo
# import streamlit as st
# from dotenv import load_dotenv
# from langchain_groq import ChatGroq
# from langchain_core.output_parsers import StrOutputParser
# from reportlab.lib.pagesizes import letter
# from reportlab.platypus import SimpleDocTemplate, Paragraph
# from reportlab.lib.styles import getSampleStyleSheet
# from pptx import Presentation
# from io import BytesIO
# from PyPDF2 import PdfReader

# load_dotenv()

# # ---------------- DATABASE ----------------
# #client = pymongo.MongoClient("mongodb://localhost:27017/")
# client = pymongo.MongoClient(os.getenv("MONGO_URI"))
# db = client["StudentDB"]
# users = db["users"]
# lessons = db["lessons"]

# # ---------------- LLM ----------------
# def LLM_Setup(prompt):
#     model = ChatGroq(
#         model="llama-3.3-70b-versatile",
#         groq_api_key=os.getenv("key")
#     )
#     parser = StrOutputParser()
#     chain = model | parser
#     return chain.invoke(prompt)

# # ---------------- SYLLABUS READER ----------------
# def read_syllabus(uploaded):
#     text=""
#     if uploaded.type=="application/pdf":
#         pdf=PdfReader(uploaded)
#         for p in pdf.pages:
#             if p.extract_text():
#                 text+=p.extract_text()
#     else:
#         text=str(uploaded.read(),"utf-8")
#     return text[:6000]

# # ---------------- PDF ----------------
# def create_pdf(text):
#     buffer=BytesIO()
#     doc=SimpleDocTemplate(buffer,pagesize=letter)
#     styles=getSampleStyleSheet()
#     content=[Paragraph(line,styles["Normal"]) for line in text.split("\n")]
#     doc.build(content)
#     buffer.seek(0)
#     return buffer

# # ---------------- PPT ----------------
# def create_ppt(text):
#     prs=Presentation()
#     blocks=text.split("\n\n")
#     for b in blocks:
#         slide=prs.slides.add_slide(prs.slide_layouts[1])
#         slide.shapes.title.text="Lesson Plan"
#         slide.placeholders[1].text=b[:1200]
#     buf=BytesIO()
#     prs.save(buf)
#     buf.seek(0)
#     return buf

# # ---------------- PAGE ----------------
# st.set_page_config(page_title="AI Lesson Planner", layout="wide")

# # ---------------- THEME ----------------
# mode = st.sidebar.radio("Theme Mode", ["Light Mode", "Dark Mode"])
# dark = mode == "Dark Mode"

# bg = "#111827" if dark else "#f5f7fb"
# text = "#f9fafb" if dark else "#111827"
# card = "#1f2937" if dark else "#ffffff"

# # ---------------- CSS ----------------
# st.markdown(f"""
# <style>

# .stApp {{background:{bg};color:{text};}}

# div.stButton > button {{
# width:240px;height:44px;display:block;margin:auto;
# border-radius:10px;font-weight:600;
# background:linear-gradient(90deg,#2563eb,#1d4ed8);
# color:white;border:none;
# }}

# div.stDownloadButton > button {{
# width:240px;height:44px;display:block;margin:auto;
# border-radius:10px;font-weight:600;
# background:linear-gradient(90deg,#2563eb,#1d4ed8)!important;
# color:white!important;border:none;
# }}

# .card {{
# background:{card};
# padding:30px;
# border-radius:14px;
# box-shadow:none;   /* SHADOW REMOVED */
# }}

# </style>
# """, unsafe_allow_html=True)

# # ---------------- HEADER ----------------
# st.markdown("<h1 style='text-align:center'>📚 AI Lesson Planner</h1>",unsafe_allow_html=True)

# # ---------------- SESSION ----------------
# if "logged_in" not in st.session_state: st.session_state.logged_in=False
# if "username" not in st.session_state: st.session_state.username=None
# if "lesson" not in st.session_state: st.session_state.lesson=""

# # =====================================================
# # LOGIN / SIGNUP
# # =====================================================
# if not st.session_state.logged_in:

#     action = st.sidebar.radio("Choose Action",["Login","Signup"])
#     c1,c2,c3 = st.columns([1,2,1])

#     with c2:
#         st.markdown("<div class='card'>",unsafe_allow_html=True)

#         if action=="Signup":
#             st.subheader("Create Account")
#             u=st.text_input("Username")
#             p=st.text_input("Password",type="password")

#             if st.button("Create Account"):
#                 if users.find_one({"username":u}):
#                     st.error("Username exists")
#                 else:
#                     users.insert_one({
#                         "username":u,
#                         "password":bcrypt.hashpw(p.encode(),bcrypt.gensalt())
#                     })
#                     st.success("Account created!")

#         else:
#             st.subheader("Login")
#             u=st.text_input("Username")
#             p=st.text_input("Password",type="password")

#             if st.button("Login"):
#                 user=users.find_one({"username":u})
#                 if user and bcrypt.checkpw(p.encode(),user["password"]):
#                     st.session_state.logged_in=True
#                     st.session_state.username=u
#                     st.rerun()
#                 else:
#                     st.error("Invalid credentials")

#         st.markdown("</div>",unsafe_allow_html=True)

# # =====================================================
# # MAIN APP
# # =====================================================
# if st.session_state.logged_in:

#     st.sidebar.success(f"👤 {st.session_state.username}")

#     page = st.sidebar.radio(
#         "📚 Teacher Menu",
#         ["Lesson Generator","Weekly Planner","Worksheet Generator",
#          "Class Activities","Quiz Generator",
#          "Upload Syllabus → Auto Plan","Saved Lessons"]
#     )

#     with st.sidebar:
#         st.header("Lesson Details")
#         subject=st.text_input("Subject")
#         topic=st.text_input("Topic")
#         grade=st.text_input("Grade")
#         duration=st.text_input("Duration")
#         difficulty=st.selectbox("Difficulty Level",["Easy","Medium","Hard"])
#         obj=st.text_area("Learning Objectives")
#         uploaded = st.file_uploader("📂 Upload Syllabus PDF/TXT", type=["pdf","txt"])

#     def center_button(label):
#         c1,c2,c3 = st.columns([1,2,1])
#         with c2:
#             return st.button(label)

#     if page=="Lesson Generator":
#         if center_button("Generate Lesson Plan"):
#             st.session_state.lesson=LLM_Setup(f"""
# Create professional lesson plan.
# Subject:{subject}
# Topic:{topic}
# Grade:{grade}
# Duration:{duration}
# Difficulty:{difficulty}
# Objectives:{obj}
# Return Markdown
# """)

#     elif page=="Weekly Planner":
#         if center_button("Generate Weekly Plan"):
#             st.session_state.lesson=LLM_Setup(
# f"Create structured 5-day weekly teaching plan for {subject} {topic} grade {grade}"
# )

#     elif page=="Worksheet Generator":
#         if center_button("Generate Worksheet"):
#             st.session_state.lesson=LLM_Setup(
# f"Create printable worksheet for {subject} {topic} grade {grade} difficulty {difficulty}"
# )

#     elif page=="Class Activities":
#         if center_button("Generate Activities"):
#             st.session_state.lesson=LLM_Setup(
# f"Create engaging classroom activities for {subject} {topic} grade {grade}"
# )

#     elif page=="Quiz Generator":
#         qnum = st.slider("Number of Questions",5,25,10)
#         if center_button("Generate Quiz"):
#             st.session_state.lesson = LLM_Setup(f"""
# Create classroom quiz.

# Subject:{subject}
# Topic:{topic}
# Grade:{grade}
# Difficulty:{difficulty}

# Generate {qnum} MCQ questions with 4 options and correct answer.

# Return Markdown
# """)

#     elif page=="Upload Syllabus → Auto Plan":

#         if uploaded is None:
#             st.info("Upload syllabus first")

#         else:
#             if center_button("Generate From Syllabus"):

#                 syllabus_text=read_syllabus(uploaded)

#                 st.session_state.lesson = LLM_Setup(f"""
# Using this syllabus:

# {syllabus_text}

# Create:

# ✔ full lesson plan
# ✔ weekly breakdown
# ✔ classroom activities
# ✔ assessment
# ✔ homework

# Return Markdown
# """)

#     if st.session_state.lesson:

#         st.success("Generated Successfully")
#         st.markdown(st.session_state.lesson)

#         if center_button("💾 Save Lesson"):
#             lessons.insert_one({
#                 "username":st.session_state.username,
#                 "lesson":st.session_state.lesson
#             })
#             st.success("Saved!")

#         pdf=create_pdf(st.session_state.lesson)
#         st.download_button("⬇ Download PDF",pdf,"lesson.pdf")

#         ppt=create_ppt(st.session_state.lesson)
#         st.download_button("⬇ Download PPT",ppt,"lesson.pptx")

# if page=="Saved Lessons":
#     st.title("Saved Lessons")
#     data=lessons.find({"username":st.session_state.username})
#     found=False
#     for d in data:
#         found=True
#         st.markdown("---")
#         st.markdown(d["lesson"])
#     if not found:
#         st.info("No saved lessons yet.")
# import os
# import bcrypt
# import pymongo
# import streamlit as st
# from dotenv import load_dotenv
# from langchain_groq import ChatGroq
# from langchain_core.output_parsers import StrOutputParser
# from reportlab.lib.pagesizes import letter
# from reportlab.platypus import SimpleDocTemplate, Paragraph
# from reportlab.lib.styles import getSampleStyleSheet
# from pptx import Presentation
# from io import BytesIO
# from PyPDF2 import PdfReader

# load_dotenv()

# # ---------------- DATABASE ----------------
# #client = pymongo.MongoClient(os.getenv("MONGO_URI"))

# try:
#     MONGO_URI = st.secrets["MONGO_URI"]
#     GROQ_KEY = st.secrets["GROQ_API_KEY"]
# except:
#     MONGO_URI = os.getenv("MONGO_URI")
#     GROQ_KEY = os.getenv("GROQ_API_KEY")

# client = pymongo.MongoClient(MONGO_URI)
# db = client["StudentDB"]
# users = db["users"]
# lessons = db["lessons"]

# # ---------------- LLM ----------------
# def LLM_Setup(prompt):
#     model = ChatGroq(
#         model="llama-3.3-70b-versatile",
#         # groq_api_key=os.getenv("GROQ_KEY")
#          groq_api_key=GROQ_KEY
#     )
#     parser = StrOutputParser()
#     chain = model | parser
#     return chain.invoke(prompt)

# # ---------------- SYLLABUS READER ----------------
# def read_syllabus(uploaded):
#     text = ""
#     if uploaded.type == "application/pdf":
#         pdf = PdfReader(uploaded)
#         for p in pdf.pages:
#             if p.extract_text():
#                 text += p.extract_text()
#     else:
#         text = str(uploaded.read(), "utf-8")
#     return text[:6000]

# # ---------------- PDF ----------------
# def create_pdf(text):
#     buffer = BytesIO()
#     doc = SimpleDocTemplate(buffer, pagesize=letter)
#     styles = getSampleStyleSheet()
#     content = [Paragraph(line, styles["Normal"]) for line in text.split("\n")]
#     doc.build(content)
#     buffer.seek(0)
#     return buffer

# # ---------------- PPT ----------------
# def create_ppt(text):
#     prs = Presentation()
#     blocks = text.split("\n\n")
#     for b in blocks:
#         slide = prs.slides.add_slide(prs.slide_layouts[1])
#         slide.shapes.title.text = "Lesson Plan"
#         slide.placeholders[1].text = b[:1200]
#     buf = BytesIO()
#     prs.save(buf)
#     buf.seek(0)
#     return buf

# # ---------------- PAGE ----------------
# st.set_page_config(page_title="AI Lesson Planner", layout="wide")

# # ---------------- THEME ----------------
# mode = st.sidebar.radio("Theme Mode", ["Light Mode", "Dark Mode"])
# dark = mode == "Dark Mode"

# bg = "#111827" if dark else "#f5f7fb"
# text = "#f9fafb" if dark else "#111827"
# card = "#1f2937" if dark else "#ffffff"

# # ---------------- CSS ----------------
# st.markdown(f"""
# <style>
# .stApp {{background:{bg};color:{text};}}
# div.stButton > button {{
#     width:240px;height:44px;display:block;margin:auto;
#     border-radius:10px;font-weight:600;
#     background:linear-gradient(90deg,#2563eb,#1d4ed8);
#     color:white;border:none;
# }}
# div.stDownloadButton > button {{
#     width:240px;height:44px;display:block;margin:auto;
#     border-radius:10px;font-weight:600;
#     background:linear-gradient(90deg,#2563eb,#1d4ed8)!important;
#     color:white!important;border:none;
# }}
# .card {{
#     background:{card};
#     padding:30px;
#     border-radius:14px;
#     box-shadow:none;
# }}
# </style>
# """, unsafe_allow_html=True)

# # ---------------- HEADER ----------------
# st.markdown("<h1 style='text-align:center'>📚 AI Lesson Planner</h1>", unsafe_allow_html=True)

# # ---------------- SESSION ----------------
# if "logged_in" not in st.session_state: st.session_state.logged_in = False
# if "username" not in st.session_state: st.session_state.username = None
# if "lesson" not in st.session_state: st.session_state.lesson = ""

# # =====================================================
# # LOGIN / SIGNUP
# # =====================================================
# if not st.session_state.logged_in:

#     action = st.sidebar.radio("Choose Action", ["Login", "Signup"])
#     c1, c2, c3 = st.columns([1, 2, 1])

#     with c2:
#         st.markdown("<div class='card'>", unsafe_allow_html=True)

#         if action == "Signup":
#             st.subheader("Create Account")
#             u = st.text_input("Username")
#             p = st.text_input("Password", type="password")

#             if st.button("Create Account"):
#                 if users.find_one({"username": u}):
#                     st.error("Username exists")
#                 else:
#                     users.insert_one({
#                         "username": u,
#                         "password": bcrypt.hashpw(p.encode(), bcrypt.gensalt())
#                     })
#                     st.success("Account created!")

#         else:
#             st.subheader("Login")
#             u = st.text_input("Username")
#             p = st.text_input("Password", type="password")

#             if st.button("Login"):
#                 user = users.find_one({"username": u})
#                 if user and bcrypt.checkpw(p.encode(), user["password"]):
#                     st.session_state.logged_in = True
#                     st.session_state.username = u
#                     st.rerun()
#                 else:
#                     st.error("Invalid credentials")

#         st.markdown("</div>", unsafe_allow_html=True)

# # =====================================================
# # MAIN APP
# # =====================================================
# if st.session_state.logged_in:

#     st.sidebar.success(f"👤 {st.session_state.username}")

#     page = st.sidebar.radio(
#         "📚 Teacher Menu",
#         ["Lesson Generator", "Weekly Planner", "Worksheet Generator",
#          "Class Activities", "Quiz Generator",
#          "Upload Syllabus → Auto Plan", "Saved Lessons"]
#     )

#     with st.sidebar:
#         st.header("Lesson Details")
#         subject = st.text_input("Subject")
#         topic = st.text_input("Topic")
#         grade = st.text_input("Grade")
#         duration = st.text_input("Duration")
#         difficulty = st.selectbox("Difficulty Level", ["Easy", "Medium", "Hard"])
#         obj = st.text_area("Learning Objectives")
#         uploaded = st.file_uploader("📂 Upload Syllabus PDF/TXT", type=["pdf", "txt"])

#     def center_button(label):
#         c1, c2, c3 = st.columns([1, 2, 1])
#         with c2:
#             return st.button(label)

#     # ---------------- PAGE LOGIC ----------------
#     if page == "Lesson Generator":
#         if center_button("Generate Lesson Plan"):
#             st.session_state.lesson = LLM_Setup(f"""
# Create professional lesson plan.
# Subject:{subject}
# Topic:{topic}
# Grade:{grade}
# Duration:{duration}
# Difficulty:{difficulty}
# Objectives:{obj}
# Return Markdown
# """)

#     elif page == "Weekly Planner":
#         if center_button("Generate Weekly Plan"):
#             st.session_state.lesson = LLM_Setup(
#                 f"Create structured 5-day weekly teaching plan for {subject} {topic} grade {grade}"
#             )

#     elif page == "Worksheet Generator":
#         if center_button("Generate Worksheet"):
#             st.session_state.lesson = LLM_Setup(
#                 f"Create printable worksheet for {subject} {topic} grade {grade} difficulty {difficulty}"
#             )

#     elif page == "Class Activities":
#         if center_button("Generate Activities"):
#             st.session_state.lesson = LLM_Setup(
#                 f"Create engaging classroom activities for {subject} {topic} grade {grade}"
#             )

#     elif page == "Quiz Generator":
#         qnum = st.slider("Number of Questions", 5, 25, 10)
#         if center_button("Generate Quiz"):
#             st.session_state.lesson = LLM_Setup(f"""
# Create classroom quiz.

# Subject:{subject}
# Topic:{topic}
# Grade:{grade}
# Difficulty:{difficulty}

# Generate {qnum} MCQ questions with 4 options and correct answer.

# Return Markdown
# """)

#     elif page == "Upload Syllabus → Auto Plan":
#         if uploaded is None:
#             st.info("Upload syllabus first")
#         else:
#             if center_button("Generate From Syllabus"):
#                 syllabus_text = read_syllabus(uploaded)
#                 st.session_state.lesson = LLM_Setup(f"""
# Using this syllabus:

# {syllabus_text}

# Create:

# ✔ full lesson plan
# ✔ weekly breakdown
# ✔ classroom activities
# ✔ assessment
# ✔ homework

# Return Markdown
# """)

#     elif page == "Saved Lessons":
#         st.title("Saved Lessons")
#         data = lessons.find({"username": st.session_state.username})
#         found = False
#         for d in data:
#             found = True
#             st.markdown("---")
#             st.markdown(d["lesson"])
#         if not found:
#             st.info("No saved lessons yet.")

#     # ---------------- DISPLAY & DOWNLOAD ----------------
#     if st.session_state.lesson and page != "Saved Lessons":
#         st.success("Generated Successfully")
#         st.markdown(st.session_state.lesson)

#         if center_button("💾 Save Lesson"):
#             lessons.insert_one({
#                 "username": st.session_state.username,
#                 "lesson": st.session_state.lesson
#             })
#             st.success("Saved!")

#         pdf = create_pdf(st.session_state.lesson)
#         st.download_button("⬇ Download PDF", pdf, "lesson.pdf")

#         ppt = create_ppt(st.session_state.lesson)
#         st.download_button("⬇ Download PPT", ppt, "lesson.pptx")
import os
import bcrypt
import pymongo
import streamlit as st
from dotenv import load_dotenv
from langchain_groq import ChatGroq
from langchain_core.output_parsers import StrOutputParser
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from io import BytesIO
from PyPDF2 import PdfReader

# ---------------- LOAD ENV ----------------
load_dotenv()

# ---------------- SAFE SECRETS HANDLING ----------------
try:
    MONGO_URI = st.secrets["MONGO_URI"]
    GROQ_KEY = st.secrets["GROQ_API_KEY"]
except:
    MONGO_URI = os.getenv("MONGO_URI")
    GROQ_KEY = os.getenv("GROQ_API_KEY")

# ---------------- DATABASE ----------------
client = pymongo.MongoClient(MONGO_URI)
db = client["StudentDB"]
users = db["users"]
lessons = db["lessons"]

# ---------------- LLM ----------------
def LLM_Setup(prompt):
    model = ChatGroq(
        model="llama-3.3-70b-versatile",
        groq_api_key=GROQ_KEY
    )
    parser = StrOutputParser()
    return (model | parser).invoke(prompt)

# ---------------- FILE READER ----------------
def read_syllabus(uploaded):
    text = ""
    if uploaded.type == "application/pdf":
        pdf = PdfReader(uploaded)
        for p in pdf.pages:
            if p.extract_text():
                text += p.extract_text()
    else:
        text = str(uploaded.read(), "utf-8")
    return text[:6000]

# ---------------- PDF ----------------
def create_pdf(text):
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    styles = getSampleStyleSheet()
    content = [Paragraph(line, styles["Normal"]) for line in text.split("\n")]
    doc.build(content)
    buffer.seek(0)
    return buffer

# ---------------- PPT ----------------
def create_ppt(text):
    prs = Presentation()
    blocks = text.split("\n\n")
    for b in blocks:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Lesson Plan"
        slide.placeholders[1].text = b[:1200]
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

# ---------------- PAGE CONFIG ----------------
st.set_page_config(page_title="AI Lesson Planner", layout="wide")

# ---------------- THEME ----------------
mode = st.sidebar.radio("Theme Mode", ["Light Mode", "Dark Mode"])
dark = mode == "Dark Mode"

bg = "#111827" if dark else "#f5f7fb"
text = "#f9fafb" if dark else "#111827"
card = "#1f2937" if dark else "#ffffff"

# ---------------- CSS FIXED ----------------
st.markdown(f"""
<style>
.stApp {{
    background: {bg};
}}

h1, h2, h3, h4, h5, h6, p, label, span {{
    color: {text} !important;
}}

input, textarea {{
    color: black !important;
}}

div.stButton > button {{
    width:240px;height:44px;margin:auto;display:block;
    border-radius:10px;font-weight:600;
    background:linear-gradient(90deg,#2563eb,#1d4ed8);
    color:white !important;border:none;
}}

div.stDownloadButton > button {{
    width:240px;height:44px;margin:auto;display:block;
    border-radius:10px;font-weight:600;
    background:linear-gradient(90deg,#2563eb,#1d4ed8)!important;
    color:white!important;border:none;
}}

.card {{
    background:{card};
    padding:30px;
    border-radius:14px;
    margin-top:20px;
}}
</style>
""", unsafe_allow_html=True)

# ---------------- HEADER ----------------
st.markdown("<h1 style='text-align:center'>📚 AI Lesson Planner</h1>", unsafe_allow_html=True)

# ---------------- SESSION ----------------
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False
if "username" not in st.session_state:
    st.session_state.username = None
if "lesson" not in st.session_state:
    st.session_state.lesson = ""

# ---------------- LOGIN / SIGNUP ----------------
if not st.session_state.logged_in:

    action = st.sidebar.radio("Choose Action", ["Login", "Signup"])
    c1, c2, c3 = st.columns([1,2,1])

    with c2:
        st.markdown("<div class='card'>", unsafe_allow_html=True)

        if action == "Signup":
            st.subheader("Create Account")
            u = st.text_input("Username")
            p = st.text_input("Password", type="password")

            if st.button("Create Account"):
                if users.find_one({"username": u}):
                    st.error("Username exists")
                else:
                    users.insert_one({
                        "username": u,
                        "password": bcrypt.hashpw(p.encode(), bcrypt.gensalt())
                    })
                    st.success("Account created!")

        else:
            st.subheader("Login")
            u = st.text_input("Username")
            p = st.text_input("Password", type="password")

            if st.button("Login"):
                user = users.find_one({"username": u})
                if user and bcrypt.checkpw(p.encode(), user["password"]):
                    st.session_state.logged_in = True
                    st.session_state.username = u
                    st.rerun()
                else:
                    st.error("Invalid credentials")

        st.markdown("</div>", unsafe_allow_html=True)

# ---------------- MAIN APP ----------------
if st.session_state.logged_in:

    st.sidebar.success(f"👤 {st.session_state.username}")

    page = st.sidebar.radio("📚 Teacher Menu", [
        "Lesson Generator","Weekly Planner","Worksheet Generator",
        "Class Activities","Quiz Generator",
        "Upload Syllabus → Auto Plan","Saved Lessons"
    ])

    with st.sidebar:
        st.header("Lesson Details")
        subject = st.text_input("Subject")
        topic = st.text_input("Topic")
        grade = st.text_input("Grade")
        duration = st.text_input("Duration")
        difficulty = st.selectbox("Difficulty", ["Easy","Medium","Hard"])
        obj = st.text_area("Objectives")
        uploaded = st.file_uploader("Upload Syllabus", type=["pdf","txt"])

    def btn(label):
        return st.button(label, use_container_width=True)

    # ---------------- GENERATION ----------------
    if page == "Lesson Generator":
        if btn("Generate Lesson Plan"):
            st.session_state.lesson = LLM_Setup(f"Create lesson plan for {subject} {topic}")

    elif page == "Weekly Planner":
        if btn("Generate Weekly Plan"):
            st.session_state.lesson = LLM_Setup(f"Weekly plan for {subject} {topic}")

    elif page == "Worksheet Generator":
        if btn("Generate Worksheet"):
            st.session_state.lesson = LLM_Setup(f"Worksheet for {subject} {topic}")

    elif page == "Class Activities":
        if btn("Generate Activities"):
            st.session_state.lesson = LLM_Setup(f"Activities for {subject} {topic}")

    elif page == "Quiz Generator":
        if btn("Generate Quiz"):
            st.session_state.lesson = LLM_Setup(f"Quiz for {subject} {topic}")

    elif page == "Upload Syllabus → Auto Plan":
        if uploaded and btn("Generate From Syllabus"):
            txt = read_syllabus(uploaded)
            st.session_state.lesson = LLM_Setup(txt)

    elif page == "Saved Lessons":
        st.title("Saved Lessons")
        for d in lessons.find({"username": st.session_state.username}):
            st.markdown("---")
            st.markdown(d["lesson"])

    # ---------------- OUTPUT ----------------
    if st.session_state.lesson and page != "Saved Lessons":
        st.success("Generated Successfully")
        st.markdown(st.session_state.lesson)

        if btn("💾 Save Lesson"):
            lessons.insert_one({
                "username": st.session_state.username,
                "lesson": st.session_state.lesson
            })

        st.download_button("⬇ PDF", create_pdf(st.session_state.lesson), "lesson.pdf")
        st.download_button("⬇ PPT", create_ppt(st.session_state.lesson), "lesson.pptx")